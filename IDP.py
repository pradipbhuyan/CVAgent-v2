# ==============================
# INTELLIGENT Resume PROCESSOR
# SP+One Drive + JD RANKING + Detailed Assessemnt + Batch Job
# ==============================

import re
import time
import zipfile
import tempfile
import hashlib
import threading
import os
from io import BytesIO
from pathlib import Path
from copy import deepcopy
import textwrap
import json
import uuid

import pandas as pd
import streamlit as st

from docx import Document as DocxDocument
from pptx import Presentation

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.document_loaders import TextLoader, PyPDFLoader

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.lib import colors

from workflow import build_graph
from core import (
    build_resume,
    send_to_concur,
    validate_document_data,
    build_confidence_map,
    classify_exception,
    extract_text_from_pdf_with_ocr_fallback,
    ocr_image_bytes_with_vlm,
    validate_resume_template,
    detect_duplicate_document,
    score_resume_against_jd,
    generate_consolidated_assessment_data,
    build_consolidated_assessment_pdf,
)

from sharepoint_connector import (
    get_cv_files_from_sharepoint,
    get_cv_files_from_onedrive,
    get_cv_files_from_sharepoint_url,
    get_cv_files_from_onedrive_url,
)

class RemoteUploadedFile:
    def __init__(self, name: str, content: bytes):
        self.name = name
        self._content = content
        self._pos = 0

    def getvalue(self):
        return self._content

    def read(self):
        return self._content

    def seek(self, pos):
        self._pos = pos

# ------------------------------
# PAGE CONFIG
# ------------------------------
st.set_page_config("IDP - Professional", layout="wide")
USERS = st.secrets.get("users", {})
MAX_BATCH_FILES = 15

BACKGROUND_JOBS = {}
BACKGROUND_JOBS_LOCK = threading.Lock()

# ------------------------------
# CACHED MODELS
# ------------------------------
@st.cache_resource
def get_llm(api_key, model):
    return ChatOpenAI(model=model, temperature=0, api_key=api_key)


@st.cache_resource
def get_embeddings(api_key):
    return OpenAIEmbeddings(api_key=api_key)

# ------------------------------
# AUTH
# ------------------------------
def validate_api_key(api_key):
    try:
        llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0,
            api_key=api_key
        )
        llm.invoke("Reply with OK")
        return True
    except Exception:
        return False


def login():
    logo_path = Path(__file__).parent / "ResumeProcessor.png"
    col1, col2, col3 = st.columns([1, 2, 1])

    with col2:
        if logo_path.exists():
            st.image(logo_path, width=300)

        st.markdown("### Sign In")
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        api_key = st.text_input("OpenAI API Key", type="password")

        if st.button("Login", use_container_width=True):
            if username not in USERS or USERS[username]["password"] != password:
                st.error("Invalid username or password")
                return

            if not api_key:
                st.error("Please enter your OpenAI API key")
                return

            with st.spinner("Validating API key..."):
                if not validate_api_key(api_key):
                    st.error("Invalid OpenAI API key")
                    return

            st.session_state["logged_in"] = True
            st.session_state["user"] = username
            st.session_state["role"] = USERS[username].get("role", "user")
            st.session_state["api_key"] = api_key
            st.rerun()

# ------------------------------
# SESSION INIT
# ------------------------------
DEFAULT_KEYS = {
    "logged_in": False,
    "user": None,
    "role": None,
    "api_key": None,
    "model_choice": "gpt-4o-mini",
    "metrics": {
        "tokens": 0,
        "input_tokens": 0,
        "output_tokens": 0,
        "cost": 0.0,
        "response_times": [],
        "calls": 0
    },
    "doc_costs": {},
    "batch_results": [],
    "exception_queue": [],
    "active_batch_index": 0,
    "batch_processed": False,
    "last_batch_signature": None,
    "show_reprocess_confirm": False,
    "pending_batch_signature": None,
    "batch_total_files": 0,
    "batch_processed_files": 0,
    "batch_current_file": None,
    "batch_file_statuses": [],
    "batch_started_at": None,
    "batch_completed_at": None,
    "batch_elapsed_seconds": 0.0,
    "current_file_started_at": None,
    "review_data": None,
    "confidence_map": None,
    "validation_result": None,
    "duplicate_info": None,
    "vectorstore": None,
    "chat_history": [],
    "suggested_questions": [],
    "current_file": None,
    "doc_type": None,
    "full_text": None,
    "auto_result": None,
    "generated_resume": None,
    "agent_events": [],
    "agent_logs": [],
    "agent_timings": {},
    "active_agent": None,
    "current_step": "Waiting",
    "progress_value": 0,
    "live_step_placeholder": None,
    "live_progress_placeholder": None,
    "live_event_placeholder": None,
    "live_pipeline_placeholder": None,
    "uploader_key": 0,
    "source_mode": "Local Upload",
    "remote_uploaded_files": [],
    "template_library": [],
    "active_template_index": None,
    "version_history": [],
    "jd_text": "",
    "jd_rankings": [],
    "detailed_assessment_data": None,
    "detailed_assessment_pdf": None,
    "job_running": False,
    "job_status": "Idle",
    "job_progress": 0,
    "job_total_files": 0,
    "job_processed_files": 0,
    "job_current_file": None,
    "job_results": [],
    "job_exception_queue": [],
    "job_output_zip": None,
    "job_output_zip_name": None,
    "job_assessment_pdf": None,
    "job_rankings": [],
    "job_notifications": [],
    "job_thread_started": False,
    "active_background_job_id": None,
}
for key, value in DEFAULT_KEYS.items():
    if key not in st.session_state:
        st.session_state[key] = value

if not st.session_state["logged_in"]:
    login()
    st.stop()

# ------------------------------
# HELPERS
# ------------------------------
def get_job_runtime_dir():
    base = Path("job_runtime")
    base.mkdir(parents=True, exist_ok=True)
    return base


def get_job_dir(job_id: str):
    job_dir = get_job_runtime_dir() / job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    return job_dir


def get_job_meta_path(job_id: str):
    return get_job_dir(job_id) / "job.json"


def write_job_json(job_id: str, data: dict):
    path = get_job_meta_path(job_id)
    serializable = deepcopy(data)

    if "source_files" in serializable:
        normalized_files = []
        for f in serializable.get("source_files", []):
            if isinstance(f, dict):
                normalized_files.append({
                    "name": f.get("name"),
                    "content_hex": f.get("content_hex"),
                })
            else:
                normalized_files.append({
                    "name": f.name,
                    "content_hex": f.getvalue().hex(),
                })
        serializable["source_files"] = normalized_files

    path.write_text(
        json.dumps(serializable, indent=2, ensure_ascii=False),
        encoding="utf-8"
    )


def read_job_json(job_id: str):
    path = get_job_meta_path(job_id)
    if not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def save_job_binary(job_id: str, filename: str, content: bytes):
    path = get_job_dir(job_id) / filename
    path.write_bytes(content)
    return str(path)


def load_job_binary(path_str: str):
    path = Path(path_str)
    if path.exists():
        return path.read_bytes()
    return None


def get_incomplete_files_for_job(job):
    completed_names = set()

    for item in job.get("results", []):
        status = item.get("status")
        file_name = item.get("file_name")
        if status in ["Completed", "Review Needed"] and file_name:
            completed_names.add(file_name)

    remaining = []
    for f in restore_source_files_from_job(job):
        if getattr(f, "name", None) not in completed_names:
            remaining.append(f)

    return remaining


def can_resume_background_job(job):
    if not job:
        return False

    status = str(job.get("status", ""))
    if job.get("is_running"):
        return False

    if status == "Completed":
        return False

    remaining = get_incomplete_files_for_job(job)
    return len(remaining) > 0


def resume_background_batch_job():
    job_id = st.session_state.get("active_background_job_id")
    if not job_id:
        st.warning("No background job found to resume.")
        return

    job = get_background_job(job_id)
    if not job:
        st.warning("Background job not found.")
        return

    if job.get("is_running"):
        st.warning("Background job is already running.")
        return

    if not can_resume_background_job(job):
        st.warning("No resumable files found for this batch job.")
        return

    jd_text = job.get("jd_text", "")
    remaining_files = get_incomplete_files_for_job(job)

    update_background_job(
        job_id,
        status="Queued for Resume",
        is_running=True,
    )

    thread = threading.Thread(
        target=run_background_batch_job,
        args=(job_id, remaining_files, jd_text),
        daemon=True,
    )
    thread.start()
    st.success(f"Background batch job resumed for {len(remaining_files)} remaining file(s).")

def create_background_job_record(uploaded_files, jd_text):
    job_id = str(uuid.uuid4())

    active_template = get_active_template_bytes()

    job = {
        "job_id": job_id,
        "status": "Queued",
        "progress": 0,
        "total_files": len(uploaded_files),
        "processed_files": 0,
        "current_file": None,
        "results": [],
        "exceptions": [],
        "output_zip_path": None,
        "output_zip_name": None,
        "assessment_pdf_path": None,
        "rankings": [],
        "notifications": [],
        "jd_text": jd_text,
        "source_files": uploaded_files,
        "template_bytes_hex": active_template.hex() if active_template else None,
        "is_running": False,
        "last_error": None,
        "last_update": time.time(),
        "created_at": time.time(),
    }

    write_job_json(job_id, job)
    return job_id


def update_background_job(job_id, **kwargs):
    job = read_job_json(job_id)
    if not job:
        return

    job.update(kwargs)
    job["last_update"] = time.time()
    write_job_json(job_id, job)


def append_background_job_notification(job_id, message):
    job = read_job_json(job_id)
    if not job:
        return

    notifications = list(job.get("notifications", []))
    notifications.append(message)
    job["notifications"] = notifications
    job["last_update"] = time.time()
    write_job_json(job_id, job)


def get_background_job(job_id):
    return read_job_json(job_id)

def reset_source_state():
    st.session_state.batch_results = []
    st.session_state.exception_queue = []
    st.session_state.active_batch_index = 0
    st.session_state.batch_processed = False
    st.session_state.last_batch_signature = None
    st.session_state.show_reprocess_confirm = False
    st.session_state.pending_batch_signature = None
    st.session_state.batch_total_files = 0
    st.session_state.batch_processed_files = 0
    st.session_state.batch_current_file = None
    st.session_state.batch_file_statuses = []
    st.session_state.jd_rankings = []
    st.session_state.jd_text = ""
    st.session_state.remote_uploaded_files = []
    st.session_state.source_mode = "Local Upload"
    reset_run_state()

def restore_source_files_from_job(job):
    files = []
    for item in job.get("source_files", []) or []:
        if isinstance(item, dict):
            name = item.get("name")
            content_hex = item.get("content_hex", "")
            files.append(
                RemoteUploadedFile(
                    name=name,
                    content=bytes.fromhex(content_hex) if content_hex else b"",
                )
            )
    return files

def extract_jd_text_from_upload(uploaded_file):
    if not uploaded_file:
        return ""

    suffix = Path(uploaded_file.name).suffix.lower()
    file_path = save_temp_file(uploaded_file)

    try:
        if suffix == ".pdf":
            docs = PyPDFLoader(file_path).load()
            return "\n".join(
                [d.page_content for d in docs if getattr(d, "page_content", None)]
            ).strip()

        if suffix == ".docx":
            return extract_docx_text(file_path).strip()

    except Exception as e:
        st.error(f"JD file read failed: {str(e)}")
        return ""

    st.warning("Unsupported JD file type. Please upload PDF or DOCX.")
    return ""




def reset_run_state():
    st.session_state["review_data"] = None
    st.session_state["confidence_map"] = None
    st.session_state["validation_result"] = None
    st.session_state["duplicate_info"] = None
    st.session_state["vectorstore"] = None
    st.session_state["chat_history"] = []
    st.session_state["suggested_questions"] = []
    st.session_state["current_file"] = None
    st.session_state["doc_type"] = None
    st.session_state["full_text"] = None
    st.session_state["auto_result"] = None
    st.session_state["generated_resume"] = None
    st.session_state["agent_events"] = []
    st.session_state["agent_logs"] = []
    st.session_state["current_step"] = "Waiting"
    st.session_state["progress_value"] = 0
    st.session_state["duplicate_info"] = None
    st.session_state["agent_timings"] = {}
    st.session_state["active_agent"] = None


def reset_single_file_state():
    st.session_state["review_data"] = None
    st.session_state["confidence_map"] = None
    st.session_state["validation_result"] = None
    st.session_state["duplicate_info"] = None
    st.session_state["vectorstore"] = None
    st.session_state["chat_history"] = []
    st.session_state["suggested_questions"] = []
    st.session_state["current_file"] = None
    st.session_state["doc_type"] = None
    st.session_state["full_text"] = None
    st.session_state["auto_result"] = None
    st.session_state["generated_resume"] = None
    st.session_state["agent_events"] = []
    st.session_state["agent_logs"] = []
    st.session_state["current_step"] = "Waiting"
    st.session_state["progress_value"] = 0
    st.session_state["agent_timings"] = {}
    st.session_state["active_agent"] = None

def reset_background_job_state():
    st.session_state["job_running"] = False
    st.session_state["job_status"] = "Idle"
    st.session_state["job_progress"] = 0
    st.session_state["job_total_files"] = 0
    st.session_state["job_processed_files"] = 0
    st.session_state["job_current_file"] = None
    st.session_state["job_results"] = []
    st.session_state["job_exception_queue"] = []
    st.session_state["job_output_zip"] = None
    st.session_state["job_output_zip_name"] = None
    st.session_state["job_assessment_pdf"] = None
    st.session_state["job_rankings"] = []

def save_temp_file(uploaded_file):
    suffix = Path(uploaded_file.name).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        return tmp.name


def load_default_resume_template_bytes():
    possible_paths = [
        Path("templates/resume_template.docx"),
        Path("templates:resume_template.docx"),
        Path(__file__).parent / "templates" / "resume_template.docx",
        Path(__file__).parent / "templates:resume_template.docx",
    ]
    for path in possible_paths:
        if path.exists():
            with open(path, "rb") as file:
                return file.read()
    return None


def get_active_template_bytes():
    library = st.session_state.get("template_library", [])
    active_index = st.session_state.get("active_template_index")

    if active_index is not None and 0 <= active_index < len(library):
        return library[active_index]["content"]

    return load_default_resume_template_bytes()


def add_template_to_library(uploaded_template):
    if not uploaded_template:
        return

    content = uploaded_template.getvalue()
    validation = validate_resume_template(content)

    entry = {
        "name": uploaded_template.name,
        "content": content,
        "validation": validation,
    }

    st.session_state.template_library.append(entry)
    st.session_state.active_template_index = len(st.session_state.template_library) - 1


def save_version_snapshot(file_name, doc_type, review_data, auto_result, status, note=""):
    snapshot = {
        "file_name": file_name,
        "doc_type": doc_type,
        "timestamp": pd.Timestamp.now().strftime("%Y-%m-%d %H:%M:%S"),
        "status": status,
        "note": note,
        "review_data": deepcopy(review_data) if review_data else {},
        "auto_result": deepcopy(auto_result) if auto_result else {},
    }
    st.session_state.version_history.append(snapshot)


def extract_docx_text(file_path):
    doc = DocxDocument(file_path)
    parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text and paragraph.text.strip():
            parts.append(paragraph.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts).strip()


def process_file_with_fallback(uploaded_file):
    suffix = Path(uploaded_file.name).suffix.lower()
    uploaded_file.seek(0)

    if suffix in [".png", ".jpg", ".jpeg"]:
        image_bytes = uploaded_file.getvalue()
        mime_type = "image/jpeg" if suffix in [".jpg", ".jpeg"] else "image/png"
        text = ocr_image_bytes_with_vlm(image_bytes, mime_type=mime_type)
        return {
            "documents": [Document(page_content=text)] if text else [],
            "text": text,
            "ocr_used": True,
            "extraction_mode": "image_vlm_ocr",
            "exception_reason": None if text else "OCR failed on image",
        }

    file_path = save_temp_file(uploaded_file)

    try:
        if suffix == ".txt":
            try:
                docs = TextLoader(file_path, encoding="utf-8").load()
            except Exception:
                docs = TextLoader(file_path, encoding="cp1252").load()

            text = "\n".join([d.page_content for d in docs]).strip()
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "plain_text",
                "exception_reason": None,
            }

        if suffix == ".pdf":
            pdf_result = extract_text_from_pdf_with_ocr_fallback(file_path)
            docs = [Document(page_content=pdf_result["text"])] if pdf_result["text"] else []
            return {
                "documents": docs,
                "text": pdf_result["text"],
                "ocr_used": pdf_result["ocr_used"],
                "extraction_mode": pdf_result["extraction_mode"],
                "exception_reason": pdf_result["exception_reason"],
            }

        if suffix == ".docx":
            text = extract_docx_text(file_path)
            docs = [Document(page_content=text)] if text else []
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "docx_text",
                "exception_reason": None if text else "No extractable text in DOCX",
            }

        if suffix == ".pptx":
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text_parts.append(shape.text.strip())
            text = "\n".join(text_parts).strip()
            docs = [Document(page_content=text)] if text else []
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "pptx_text",
                "exception_reason": None if text else "No extractable text in PPTX",
            }

        if suffix == ".xlsx":
            excel_file = pd.ExcelFile(file_path)
            sheet_texts = []
            for sheet in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                sheet_texts.append(f"Sheet: {sheet}")
                sheet_texts.append(df.to_string(index=False))
            text = "\n\n".join(sheet_texts).strip()
            docs = [Document(page_content=text)] if text else []
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "xlsx_text",
                "exception_reason": None if text else "No extractable text in Excel",
            }

    except Exception as e:
        return {
            "documents": [],
            "text": "",
            "ocr_used": False,
            "extraction_mode": "failed",
            "exception_reason": str(e),
        }

    return {
        "documents": [],
        "text": "",
        "ocr_used": False,
        "extraction_mode": "unsupported",
        "exception_reason": f"Unsupported file type: {suffix}",
    }


def create_vectorstore(docs):
    if not docs:
        return None

    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    chunks = splitter.split_documents(docs)
    if not chunks:
        return None

    for chunk in chunks:
        chunk.metadata = {"source": st.session_state.get("current_file", "unknown")}

    try:
        emb = get_embeddings(st.session_state["api_key"])
        return Chroma.from_documents(chunks, embedding=emb)
    except Exception:
        return None


def push_agent_log(message):
    st.session_state.agent_logs.append(message)
    refresh_live_batch_activity()


def record_agent_event(step, status, message=""):
    now = time.time()

    if "agent_timings" not in st.session_state:
        st.session_state["agent_timings"] = {}

    if status == "running":
        if step not in st.session_state["agent_timings"]:
            st.session_state["agent_timings"][step] = {}
        if not st.session_state["agent_timings"][step].get("started_at"):
            st.session_state["agent_timings"][step]["started_at"] = now
        st.session_state["active_agent"] = step

    elif status in ["done", "error"]:
        if step not in st.session_state["agent_timings"]:
            st.session_state["agent_timings"][step] = {}
        started_at = st.session_state["agent_timings"][step].get("started_at")
        st.session_state["agent_timings"][step]["ended_at"] = now
        if started_at:
            st.session_state["agent_timings"][step]["elapsed"] = round(now - started_at, 2)
        if st.session_state.get("active_agent") == step:
            st.session_state["active_agent"] = None

    st.session_state.agent_events.append({
        "step": step,
        "status": status,
        "message": message,
    })
    refresh_live_batch_activity()
    
def refresh_live_batch_activity():
    step_placeholder = st.session_state.get("live_step_placeholder")
    progress_placeholder = st.session_state.get("live_progress_placeholder")
    event_placeholder = st.session_state.get("live_event_placeholder")

    total_files = st.session_state.get("batch_total_files", 0)
    processed_files = st.session_state.get("batch_processed_files", 0)
    current_file = st.session_state.get("batch_current_file")
    current_step = st.session_state.get("current_step", "Waiting")
    file_statuses = st.session_state.get("batch_file_statuses", [])
    exception_count = len(st.session_state.get("exception_queue", []))
    per_file_progress = int(st.session_state.get("progress_value", 0))

    if total_files > 0:
        overall_progress = int(((processed_files + (per_file_progress / 100.0)) / total_files) * 100)
        overall_progress = max(0, min(100, overall_progress))
    else:
        overall_progress = per_file_progress

    if step_placeholder is not None:
        elapsed = st.session_state.get("batch_elapsed_seconds", 0.0)

        if total_files > 0:
            elapsed_line = f"**Elapsed:** {elapsed:.2f} sec  " if elapsed > 0 else ""

            step_placeholder.markdown(
                f"""
#### Batch Progress

**Current File:** {current_file or "-"}  
**Current Step:** {current_step}  
**Processed:** {processed_files} / {total_files}  
**Exceptions:** {exception_count}  
{elapsed_line}
"""
            )
        else:
            if current_step != "Waiting":
                step_placeholder.markdown(f"#### Progress\n\n**Current Step:** {current_step}")
            else:
                step_placeholder.empty()

    if progress_placeholder is not None:
        if total_files > 0 or per_file_progress > 0:
            progress_placeholder.progress(overall_progress)
        else:
            progress_placeholder.empty()

    if event_placeholder is not None:
        content = []

        if total_files > 0:
            content.append("#### File Queue")

            if file_statuses:
                for item in file_statuses:
                    status = item.get("status", "pending")
                    file_name = item.get("file_name", "")

                    if status == "done":
                        icon = "✅"
                    elif status == "error":
                        icon = "❌"
                    elif status == "running":
                        icon = "🔄"
                    else:
                        icon = "⏳"

                    line = f"{icon} **{file_name}**"
                    if item.get("message"):
                        line += f"  \n{item.get('message')}"
                    content.append(line)
            else:
                content.append("_No files started yet_")
        else:
            events = st.session_state.get("agent_events", [])
            if events:
                content.append("#### Completed Steps")
                for event in events[-8:]:
                    status = event.get("status", "pending")
                    if status == "done":
                        icon = "✅"
                    elif status == "error":
                        icon = "❌"
                    elif status == "running":
                        icon = "🔄"
                    else:
                        icon = "⏳"

                    line = f"{icon} **{event.get('step', '')}**"
                    if event.get("message"):
                        line += f"  \n{event.get('message')}"
                    content.append(line)

        event_placeholder.markdown("\n\n".join(content) if content else "")

    render_agent_pipeline()



def render_agent_pipeline():
    pipeline_placeholder = st.session_state.get("live_pipeline_placeholder")
    if pipeline_placeholder is None:
        return

    doc_type = st.session_state.get("doc_type")
    events = st.session_state.get("agent_events", [])
    timings = st.session_state.get("agent_timings", {})
    active_agent = st.session_state.get("active_agent")

    pipeline = [
        "Ingestion Agent",
        "Extraction Agent",
        "Retrieval Agent",
        "Classification Agent",
        "Structuring Agent",
        "Validation Agent",
        "Output Agent",
    ]

    if doc_type in ["invoice", "ticket"]:
        pipeline.append("Concur Agent")

    status_map = {name: {"status": "pending", "message": ""} for name in pipeline}

    for event in events:
        step = event.get("step")
        if step in status_map:
            status_map[step] = {
                "status": event.get("status", "pending"),
                "message": event.get("message", ""),
            }

    html_parts = [
        textwrap.dedent("""
        <div style="margin-top:10px;">
            <div style="font-weight:700;font-size:16px;margin-bottom:10px;">
                Agentic Pipeline
            </div>
            <div style="display:flex;flex-wrap:wrap;gap:10px;">
        """).strip()
    ]

    for step in pipeline:
        item = status_map[step]
        status = item["status"]
        elapsed = timings.get(step, {}).get("elapsed")
        running_since = timings.get(step, {}).get("started_at")

        short_name = step.replace(" Agent", "")

        if status == "done":
            bg = "#e8f7ee"
            border = "#8fd19e"
            icon = "✅"
            text = "#166534"
        elif status == "running" or step == active_agent:
            bg = "#eef4ff"
            border = "#7aa2ff"
            icon = "🔄"
            text = "#1d4ed8"
        elif status == "error":
            bg = "#fdecec"
            border = "#f5a3a3"
            icon = "❌"
            text = "#b42318"
        else:
            bg = "#f5f5f5"
            border = "#dddddd"
            icon = "⏳"
            text = "#6b7280"

        if elapsed is not None:
            subtitle = f"{elapsed:.2f}s"
        elif (status == "running" or step == active_agent) and running_since:
            subtitle = f"{round(time.time() - running_since, 2)}s"
        else:
            subtitle = item.get("message") or "Pending"

        card_html = textwrap.dedent(f"""
        <div style="
            min-width:120px;
            flex:1 1 120px;
            padding:12px 10px;
            border-radius:14px;
            border:1px solid {border};
            background:{bg};
            text-align:center;
        ">
            <div style="font-size:18px;line-height:1;">{icon}</div>
            <div style="font-weight:700;color:{text};font-size:12px;margin-top:6px;">
                {short_name}
            </div>
            <div style="font-size:11px;color:#4b5563;margin-top:4px;">
                {subtitle}
            </div>
        </div>
        """).strip()

        html_parts.append(card_html)

    html_parts.append("</div></div>")

    pipeline_placeholder.markdown("".join(html_parts), unsafe_allow_html=True)


def update_batch_file_status(file_name, status, message=""):
    statuses = st.session_state.get("batch_file_statuses", [])

    found = False
    for item in statuses:
        if item.get("file_name") == file_name:
            item["status"] = status
            item["message"] = message
            found = True
            break

    if not found:
        statuses.append({
            "file_name": file_name,
            "status": status,
            "message": message
        })

    st.session_state["batch_file_statuses"] = statuses
    refresh_live_batch_activity()


def update_progress(percent, message):
    st.session_state["progress_value"] = percent
    st.session_state["current_step"] = message

    current_file = st.session_state.get("batch_current_file")
    if current_file:
        update_batch_file_status(current_file, "running", message)

    refresh_live_batch_activity()


def get_suggested_questions(doc_type):
    if doc_type == "invoice":
        return ["What is the total amount?", "Who is the vendor?", "What is the invoice date?"]
    if doc_type == "resume":
        return ["Summarize this candidate", "What skills does the candidate have?", "What is the experience?"]
    if doc_type == "ticket":
        return ["What is the ticket number?", "What is the travel date?", "What are the key details?"]
    return ["What is this document?", "What are the key points?"]


def normalize_graph_result(result):
    if not isinstance(result, dict):
        return {
            "doc_type": None,
            "structured_data": None,
            "result": {},
            "error": "Graph returned non-dict output",
        }

    doc_type = result.get("doc_type") or result.get("type")
    structured_data = result.get("data") if doc_type in ["invoice", "ticket"] else None
    inner = result.get("result", {}) if isinstance(result.get("result", {}), dict) else {}

    return {
        "doc_type": doc_type,
        "structured_data": structured_data,
        "result": inner,
        "error": result.get("error"),
        "step_metrics": result.get("step_metrics", []),
        "confidence": result.get("confidence"),
        "validation": result.get("validation"),
        "ocr_used": result.get("ocr_used", False),
        "extraction_mode": result.get("extraction_mode"),
        "exception_reason": result.get("exception_reason"),
        "needs_review": result.get("needs_review", False),
    }

def process_single_file(uploaded_file):
    reset_single_file_state()
    st.session_state.current_file = uploaded_file.name

    record_agent_event("Ingestion Agent", "running", "Receiving file")
    update_progress(5, "Ingestion Agent — file received")
    record_agent_event("Ingestion Agent", "done", "File received")

    record_agent_event("Extraction Agent", "running", "Extracting text")
    extracted = process_file_with_fallback(uploaded_file)
    docs = extracted["documents"]
    full_text = extracted["text"]

    if extracted["ocr_used"]:
        record_agent_event("Extraction Agent", "done", "Text extracted using OCR fallback")
    else:
        record_agent_event("Extraction Agent", "done", "Text extracted")

    update_progress(20, "Extraction Agent — text extracted")

    if not full_text:
        reason = extracted["exception_reason"] or "No extractable text"
        record_agent_event("Retrieval Agent", "error", "Skipped due to missing text")
        record_agent_event("Classification Agent", "error", "Skipped due to missing text")
        record_agent_event("Structuring Agent", "error", "Skipped due to missing text")
        record_agent_event("Validation Agent", "error", "Skipped due to missing text")
        record_agent_event("Output Agent", "error", "Skipped due to missing text")
        return {
            "file_name": uploaded_file.name,
            "status": "Exception",
            "doc_type": "unknown",
            "ocr_used": extracted["ocr_used"],
            "exception_reason": reason,
            "cost": 0.0,
            "tokens": 0,
            "duplicate_info": {
                "is_duplicate": False,
                "match_file": None,
                "reason": None,
                "score": 0.0,
            },
            "agent_events": st.session_state.get("agent_events", []),
            "agent_timings": st.session_state.get("agent_timings", {}),
        }

    st.session_state.full_text = full_text

    record_agent_event("Retrieval Agent", "running", "Creating vector index")
    vectorstore = create_vectorstore(docs)
    st.session_state.vectorstore = vectorstore
    record_agent_event("Retrieval Agent", "done", "Vector index created")
    update_progress(30, "Retrieval Agent — search index ready")

    graph = build_graph()
    graph_input = {
        "text": full_text,
        "filename": uploaded_file.name,
        "template": get_active_template_bytes(),
        "progress": update_progress,
        "event_callback": record_agent_event,
        "ocr_used": extracted["ocr_used"],
        "extraction_mode": extracted["extraction_mode"],
        "exception_reason": extracted["exception_reason"],
    }

    before_cost = st.session_state["metrics"]["cost"]
    before_tokens = st.session_state["metrics"]["tokens"]

    raw_result = graph.invoke(graph_input)
    normalized = normalize_graph_result(raw_result)

    doc_type = normalized.get("doc_type")
    result = normalized.get("result", {})
    review_data = result.get("data") or normalized.get("structured_data") or {}
    
    # 🚨 HARD FILTER: ONLY allow resumes
    if doc_type != "resume":
        record_agent_event(
            "Classification Agent",
            "error",
            f"Non-CV detected: {doc_type or 'unknown'}"
        )
    
        return {
            "file_name": uploaded_file.name,
            "status": "Exception",
            "doc_type": doc_type or "unknown",
            "ocr_used": extracted["ocr_used"],
            "exception_reason": f"Not a CV/Resume (detected: {doc_type or 'unknown'})",
    
            # ⚠️ IMPORTANT: DO NOT PASS ANY RESULT DATA
            "review_data": None,
            "validation": None,
            "confidence": None,
            "duplicate_info": None,
            "auto_result": None,
            "vectorstore": None,
            "full_text": None,
    
            "cost": 0.0,
            "tokens": 0,
        }
    
    record_agent_event("Validation Agent", "running", "Checking required fields")
    validation = normalized.get("validation") or validate_document_data(review_data, doc_type)
    confidence = normalized.get("confidence") or build_confidence_map(review_data, doc_type)

    if validation.get("passed", True):
        if doc_type == "resume":
            record_agent_event("Validation Agent", "done", "Resume validation complete")
        elif doc_type == "invoice":
            record_agent_event("Validation Agent", "done", "Invoice validation complete")
        elif doc_type == "ticket":
            record_agent_event("Validation Agent", "done", "Ticket validation complete")
        else:
            record_agent_event("Validation Agent", "done", "Validation complete")
    else:
        record_agent_event("Validation Agent", "error", "Validation issues found")

    duplicate_info = detect_duplicate_document(
        new_doc_type=doc_type,
        new_data=review_data,
        existing_results=st.session_state.get("batch_results", []),
    )

    exception_reason = classify_exception(
        doc_type=doc_type,
        text=full_text,
        validation=validation,
        confidence=confidence,
        extraction_meta=extracted,
    )

    st.session_state.doc_type = doc_type
    st.session_state.review_data = review_data
    st.session_state.validation_result = validation
    st.session_state.confidence_map = confidence
    st.session_state.duplicate_info = duplicate_info
    st.session_state.auto_result = {
        "doc_type": doc_type,
        "structured_data": normalized.get("structured_data"),
        "result": result,
        "metrics": {},
        "step_metrics": normalized.get("step_metrics", []),
        "ocr_used": extracted["ocr_used"],
        "extraction_mode": extracted["extraction_mode"],
    }
    st.session_state.generated_resume = result.get("file")
    st.session_state.suggested_questions = get_suggested_questions(doc_type)

    after_cost = st.session_state["metrics"]["cost"]
    after_tokens = st.session_state["metrics"]["tokens"]

    status = "Completed"
    if exception_reason:
        status = "Exception"
    elif not validation.get("passed", True):
        status = "Review Needed"

    update_progress(100, "Workflow Agent — completed")

    save_version_snapshot(
        file_name=uploaded_file.name,
        doc_type=doc_type,
        review_data=review_data,
        auto_result=st.session_state.get("auto_result"),
        status=status,
        note="Initial extraction result"
    )

    return {
        "file_name": uploaded_file.name,
        "status": status,
        "doc_type": doc_type,
        "ocr_used": extracted["ocr_used"],
        "exception_reason": exception_reason,
        "review_data": review_data,
        "validation": validation,
        "confidence": confidence,
        "duplicate_info": duplicate_info,
        "auto_result": st.session_state.auto_result,
        "vectorstore": vectorstore,
        "full_text": full_text,
        "cost": round(after_cost - before_cost, 6),
        "tokens": after_tokens - before_tokens,
        "agent_events": deepcopy(st.session_state.get("agent_events", [])),
        "agent_timings": deepcopy(st.session_state.get("agent_timings", {})),
    }


def process_single_file_for_job(uploaded_file, existing_results=None, template_bytes=None):
    existing_results = existing_results or []

    extracted = process_file_with_fallback(uploaded_file)
    docs = extracted["documents"]
    full_text = extracted["text"]

    if not full_text:
        reason = extracted["exception_reason"] or "No extractable text"
        return {
            "file_name": uploaded_file.name,
            "status": "Exception",
            "doc_type": "unknown",
            "ocr_used": extracted["ocr_used"],
            "exception_reason": reason,
            "cost": 0.0,
            "tokens": 0,
            "duplicate_info": {
                "is_duplicate": False,
                "match_file": None,
                "reason": None,
                "score": 0.0,
            },
            "review_data": None,
            "validation": None,
            "confidence": None,
            "auto_result": None,
            "vectorstore": None,
            "full_text": None,
            "debug_info": {
                "graph_error": None,
                "detected_doc_type": None,
                "final_doc_type": "unknown",
                "resume_fallback_used": False,
                "extraction_mode": extracted.get("extraction_mode"),
                "text_preview": "",
            },
        }

    vectorstore = create_vectorstore(docs)

    graph = build_graph()
    graph_input = {
        "text": full_text,
        "filename": uploaded_file.name,
        "template": template_bytes,
        "progress": lambda percent, message: None,
        "event_callback": lambda step, status, message="": None,
        "ocr_used": extracted["ocr_used"],
        "extraction_mode": extracted["extraction_mode"],
        "exception_reason": extracted["exception_reason"],
    }

    before_cost = st.session_state.get("metrics", {}).get("cost", 0.0)
    before_tokens = st.session_state.get("metrics", {}).get("tokens", 0)

    raw_result = graph.invoke(graph_input)
    normalized = normalize_graph_result(raw_result)

    original_doc_type = normalized.get("doc_type")
    doc_type = original_doc_type
    result = normalized.get("result", {})
    review_data = result.get("data") or normalized.get("structured_data") or {}
    raw_error = normalized.get("error")
    extraction_mode = extracted.get("extraction_mode")

    resume_fallback_used = False
    if doc_type != "resume" and looks_like_resume_text(full_text):
        doc_type = "resume"
        resume_fallback_used = True

    if doc_type != "resume":
        return {
            "file_name": uploaded_file.name,
            "status": "Exception",
            "doc_type": doc_type or "unknown",
            "ocr_used": extracted["ocr_used"],
            "exception_reason": f"Not a CV/Resume (detected: {original_doc_type or 'unknown'})",
            "review_data": None,
            "validation": None,
            "confidence": None,
            "duplicate_info": None,
            "auto_result": None,
            "vectorstore": None,
            "full_text": full_text,
            "cost": 0.0,
            "tokens": 0,
            "debug_info": {
                "detected_doc_type": original_doc_type,
                "final_doc_type": doc_type,
                "graph_error": raw_error,
                "resume_fallback_used": resume_fallback_used,
                "extraction_mode": extraction_mode,
                "text_preview": full_text[:1000] if full_text else "",
            },
        }

    validation = normalized.get("validation") or validate_document_data(review_data, doc_type)
    confidence = normalized.get("confidence") or build_confidence_map(review_data, doc_type)

    duplicate_info = detect_duplicate_document(
        new_doc_type=doc_type,
        new_data=review_data,
        existing_results=existing_results,
    )

    exception_reason = classify_exception(
        doc_type=doc_type,
        text=full_text,
        validation=validation,
        confidence=confidence,
        extraction_meta=extracted,
    )

    after_cost = st.session_state.get("metrics", {}).get("cost", 0.0)
    after_tokens = st.session_state.get("metrics", {}).get("tokens", 0)

    status = "Completed"

    if exception_reason:
        reason_text = str(exception_reason).lower()

        if doc_type == "resume" and (
            "validation" in reason_text
            or "missing" in reason_text
            or "confidence" in reason_text
        ):
            status = "Review Needed"
        else:
            status = "Exception"
    elif not validation.get("passed", True):
        status = "Review Needed"

    return {
        "file_name": uploaded_file.name,
        "status": status,
        "doc_type": doc_type,
        "ocr_used": extracted["ocr_used"],
        "exception_reason": exception_reason,
        "review_data": review_data,
        "validation": validation,
        "confidence": confidence,
        "duplicate_info": duplicate_info,
        "auto_result": {
            "doc_type": doc_type,
            "structured_data": normalized.get("structured_data"),
            "result": result,
            "metrics": {},
            "step_metrics": normalized.get("step_metrics", []),
            "ocr_used": extracted["ocr_used"],
            "extraction_mode": extracted["extraction_mode"],
        },
        "vectorstore": vectorstore,
        "full_text": full_text,
        "cost": round(after_cost - before_cost, 6),
        "tokens": after_tokens - before_tokens,
        "debug_info": {
            "detected_doc_type": original_doc_type,
            "final_doc_type": doc_type,
            "graph_error": raw_error,
            "resume_fallback_used": resume_fallback_used,
            "extraction_mode": extraction_mode,
            "text_preview": full_text[:1000] if full_text else "",
        },
    }


def looks_like_resume_text(text: str) -> bool:
    if not text:
        return False

    t = text.lower()
    signals = 0

    if "@" in t or "email" in t or "email id" in t:
        signals += 1
    if "mobile" in t or "phone" in t or "contact" in t:
        signals += 1
    if "location" in t or "address" in t:
        signals += 1

    resume_keywords = [
        "experience",
        "work experience",
        "professional summary",
        "summary",
        "skills",
        "technical skills",
        "education",
        "certification",
        "projects",
        "employment",
        "career objective",
        "profile",
        "responsibilities",
    ]
    for kw in resume_keywords:
        if kw in t:
            signals += 1

    if "years of experience" in t or "year of experience" in t:
        signals += 2
    if "curriculum vitae" in t or "resume" in t:
        signals += 2

    bullet_count = text.count("•") + text.count("- ")
    if bullet_count >= 3:
        signals += 1

    return signals >= 4

def load_batch_result_into_session(index):
    if index < 0 or index >= len(st.session_state.batch_results):
        return

    item = st.session_state.batch_results[index]
    st.session_state.active_batch_index = index
    st.session_state.current_file = item.get("file_name")
    st.session_state.doc_type = item.get("doc_type")
    st.session_state.review_data = item.get("review_data")
    st.session_state.validation_result = item.get("validation")
    st.session_state.confidence_map = item.get("confidence")
    st.session_state.duplicate_info = item.get("duplicate_info")
    st.session_state.auto_result = item.get("auto_result")
    st.session_state.vectorstore = item.get("vectorstore")
    st.session_state.full_text = item.get("full_text")
    st.session_state.generated_resume = ((item.get("auto_result") or {}).get("result") or {}).get("file")
    st.session_state["agent_events"] = deepcopy(item.get("agent_events", []))
    st.session_state["agent_timings"] = deepcopy(item.get("agent_timings", {}))
    st.session_state["active_agent"] = None
    refresh_live_batch_activity()


def get_batch_signature(uploaded_files):
    if not uploaded_files:
        return None

    parts = []
    for file in uploaded_files:
        try:
            content_hash = hashlib.md5(file.getvalue()).hexdigest()
        except Exception:
            content_hash = f"{file.name}-{len(file.getvalue())}"
        parts.append(f"{file.name}:{content_hash}")

    return "|".join(parts)


def go_to_next_batch_result():
    batch_results = st.session_state.get("batch_results", [])
    if not batch_results:
        return

    current_index = st.session_state.get("active_batch_index", 0)
    next_index = current_index + 1

    if next_index < len(batch_results):
        load_batch_result_into_session(next_index)


def build_zip_from_batch_results(target_type: str) -> bytes:
    output = BytesIO()

    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zf:
        for item in st.session_state.get("batch_results", []):
            auto_result = item.get("auto_result") or {}
            result = auto_result.get("result") or {}
            doc_type = item.get("doc_type")

            if target_type == "resume" and doc_type == "resume":
                file_bytes = result.get("file")
                file_name = result.get("file_name") or f"{item.get('file_name', 'resume')}.docx"
                if file_bytes:
                    if not file_name.lower().endswith(".docx"):
                        file_name = f"{file_name}.docx"
                    zf.writestr(file_name, file_bytes)

            elif target_type == "invoice" and doc_type == "invoice":
                excel_bytes = result.get("excel")
                review_data = item.get("review_data") or {}
                file_name = (
                    review_data.get("invoice_number")
                    or review_data.get("invoice_no")
                    or review_data.get("vendor")
                    or item.get("file_name")
                    or "invoice_data"
                )
                file_name = str(file_name).strip()
                file_name = re.sub(r'[\\/*?:"<>|]', "", file_name)

                if excel_bytes:
                    if not file_name.lower().endswith(".xlsx"):
                        file_name = f"{file_name}.xlsx"
                    zf.writestr(file_name, excel_bytes)

    output.seek(0)
    return output.getvalue()

def build_zip_from_results(results, target_type: str) -> bytes:
    output = BytesIO()

    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zf:
        for item in results:
            auto_result = item.get("auto_result") or {}
            result = auto_result.get("result") or {}
            doc_type = item.get("doc_type")

            if target_type == "resume" and doc_type == "resume":
                file_bytes = result.get("file")
                file_name = result.get("file_name") or f"{item.get('file_name', 'resume')}.docx"
                if file_bytes:
                    if not file_name.lower().endswith(".docx"):
                        file_name = f"{file_name}.docx"
                    zf.writestr(file_name, file_bytes)

    output.seek(0)
    return output.getvalue()

def get_batch_download_counts():
    resume_count = 0
    invoice_count = 0

    for item in st.session_state.get("batch_results", []):
        auto_result = item.get("auto_result") or {}
        result = auto_result.get("result") or {}
        doc_type = item.get("doc_type")

        if doc_type == "resume" and result.get("file"):
            resume_count += 1

        if doc_type == "invoice" and result.get("excel"):
            invoice_count += 1

    return resume_count, invoice_count


def rank_all_resumes_against_jd():
    jd_text = (st.session_state.get("jd_text") or "").strip()
    if not jd_text:
        st.warning("Please provide a JD first.")
        return

    resume_items = [
        item for item in st.session_state.get("batch_results", [])
        if item.get("doc_type") == "resume" and item.get("review_data")
    ]

    if not resume_items:
        st.warning("No processed resumes found in the current batch.")
        return

    rankings = []
    for item in resume_items:
        resume_data = item.get("review_data") or {}
        score = score_resume_against_jd(resume_data, jd_text)
        score["file_name"] = item.get("file_name")
        rankings.append(score)

    rankings = sorted(rankings, key=lambda x: x.get("overall_score", 0), reverse=True)
    for idx, row in enumerate(rankings, start=1):
        row["rank"] = idx

    st.session_state.jd_rankings = rankings

def rank_resumes_against_jd_for_results(batch_results, jd_text):
    resume_items = [
        item for item in batch_results
        if item.get("doc_type") == "resume" and item.get("review_data")
    ]

    rankings = []
    for item in resume_items:
        resume_data = item.get("review_data") or {}
        score = score_resume_against_jd(resume_data, jd_text)
        score["file_name"] = item.get("file_name")
        rankings.append(score)

    rankings = sorted(rankings, key=lambda x: x.get("overall_score", 0), reverse=True)
    for idx, row in enumerate(rankings, start=1):
        row["rank"] = idx

    return rankings

def compact_field(label, value):
    st.markdown(
        f"**{label}**  \n<small>{value if value not in [None, ''] else '-'}</small>",
        unsafe_allow_html=True
    )



def run_background_batch_job(job_id, uploaded_files, jd_text):
    try:
        update_background_job(
            job_id,
            status="Running",
            progress=0,
            processed_files=0,
            current_file=None,
            results=[],
            exceptions=[],
            output_zip_path=None,
            output_zip_name=None,
            assessment_pdf_path=None,
            rankings=[],
            is_running=True,
            last_error=None,
        )

        local_results = []
        local_exceptions = []
        total_files = len(uploaded_files)

        for idx, uploaded_file in enumerate(uploaded_files, start=1):
            update_background_job(
                job_id,
                current_file=uploaded_file.name,
                status=f"Processing {uploaded_file.name}",
                progress=int(((idx - 1) / max(total_files, 1)) * 100),
                results=local_results,
                exceptions=local_exceptions,
            )

            saved_state = {
                "current_file": st.session_state.get("current_file"),
                "batch_results": deepcopy(st.session_state.get("batch_results", [])),
                "exception_queue": deepcopy(st.session_state.get("exception_queue", [])),
                "review_data": deepcopy(st.session_state.get("review_data")),
                "confidence_map": deepcopy(st.session_state.get("confidence_map")),
                "validation_result": deepcopy(st.session_state.get("validation_result")),
                "duplicate_info": deepcopy(st.session_state.get("duplicate_info")),
                "vectorstore": st.session_state.get("vectorstore"),
                "chat_history": deepcopy(st.session_state.get("chat_history", [])),
                "suggested_questions": deepcopy(st.session_state.get("suggested_questions", [])),
                "doc_type": st.session_state.get("doc_type"),
                "full_text": st.session_state.get("full_text"),
                "auto_result": deepcopy(st.session_state.get("auto_result")),
                "generated_resume": st.session_state.get("generated_resume"),
                "agent_events": deepcopy(st.session_state.get("agent_events", [])),
                "agent_logs": deepcopy(st.session_state.get("agent_logs", [])),
                "agent_timings": deepcopy(st.session_state.get("agent_timings", {})),
                "active_agent": st.session_state.get("active_agent"),
                "current_step": st.session_state.get("current_step", "Waiting"),
                "progress_value": st.session_state.get("progress_value", 0),
            }

            try:
                st.session_state["batch_results"] = local_results
                st.session_state["exception_queue"] = local_exceptions

                result = process_single_file(uploaded_file)
                local_results.append(result)

                if result.get("status") == "Exception":
                    local_exceptions.append(result)

                update_background_job(
                    job_id,
                    status=f"Finished {uploaded_file.name}",
                    results=local_results,
                    exceptions=local_exceptions,
                )

            except Exception as e:
                error_result = {
                    "file_name": uploaded_file.name,
                    "status": "Exception",
                    "doc_type": "unknown",
                    "ocr_used": False,
                    "exception_reason": f"Unhandled error: {str(e)}",
                    "cost": 0.0,
                    "tokens": 0,
                    "duplicate_info": {
                        "is_duplicate": False,
                        "match_file": None,
                        "reason": None,
                        "score": 0.0,
                    },
                    "review_data": None,
                    "validation": None,
                    "confidence": None,
                    "auto_result": None,
                    "vectorstore": None,
                    "full_text": None,
                }
                local_results.append(error_result)
                local_exceptions.append(error_result)

                update_background_job(
                    job_id,
                    status=f"Error in {uploaded_file.name}",
                    results=local_results,
                    exceptions=local_exceptions,
                    last_error=str(e),
                )

            finally:
                for k, v in saved_state.items():
                    st.session_state[k] = v

            update_background_job(
                job_id,
                processed_files=idx,
                progress=int((idx / max(total_files, 1)) * 100),
            )

        rankings = rank_resumes_against_jd_for_results(local_results, jd_text)
        assessment_pdf_path = None

        if rankings:
            update_background_job(job_id, status="Generating assessment report")
            report_data = generate_consolidated_assessment_data(
                batch_results=local_results,
                jd_text=jd_text,
                jd_rankings=rankings,
            )
            pdf_bytes = build_consolidated_assessment_pdf(report_data)
            assessment_pdf_path = save_job_binary(job_id, "BackgroundDetailedAssessment.pdf", pdf_bytes)

        update_background_job(job_id, status="Building resume ZIP")
        resume_zip = build_zip_from_results(local_results, "resume")
        zip_path = save_job_binary(job_id, "background_job_resumes.zip", resume_zip)

        review_needed_count = len([x for x in local_results if x.get("status") == "Review Needed"])
        exception_count = len([x for x in local_results if x.get("status") == "Exception"])

        if exception_count > 0:
            final_status = "Completed with Exceptions"
        elif review_needed_count > 0:
            final_status = "Completed with Review Needed"
        else:
            final_status = "Completed"

        update_background_job(
            job_id,
            status=final_status,
            progress=100,
            current_file=None,
            results=local_results,
            exceptions=local_exceptions,
            rankings=rankings,
            assessment_pdf_path=assessment_pdf_path,
            output_zip_path=zip_path,
            output_zip_name="background_job_resumes.zip",
            is_running=False,
        )

        append_background_job_notification(
            job_id,
            f"Background batch completed successfully. Processed {len(local_results)} file(s)."
        )

    except Exception as e:
        update_background_job(
            job_id,
            status=f"Failed: {str(e)}",
            current_file=None,
            is_running=False,
            last_error=str(e),
        )
        append_background_job_notification(
            job_id,
            f"Background batch failed: {str(e)}"
        )


def submit_background_batch_job(uploaded_files):
    jd_text = (st.session_state.get("jd_text") or "").strip()

    if not jd_text:
        st.warning("Background batch job cannot start because JD is missing.")
        return

    if not uploaded_files:
        st.warning("No CV files available for background job.")
        return

    active_job_id = st.session_state.get("active_background_job_id")
    if active_job_id:
        active_job = get_background_job(active_job_id)
        if active_job and active_job.get("is_running"):
            st.warning("A background batch job is already running.")
            return

    job_id = create_background_job_record(uploaded_files, jd_text)
    st.session_state["active_background_job_id"] = job_id

    thread = threading.Thread(
        target=run_background_batch_job,
        args=(job_id, uploaded_files, jd_text),
        daemon=True,
    )
    thread.start()
    st.success("Background batch job started.")


def render_job_notifications():
    job_id = st.session_state.get("active_background_job_id")
    if not job_id:
        return

    job = get_background_job(job_id)
    if not job:
        return

    notifications = job.get("notifications", [])
    if notifications:
        latest = notifications[-1]
        try:
            st.toast(latest)
        except Exception:
            st.success(latest)

def render_background_job_monitor():
    st.markdown("### Background Batch Job")

    job_id = st.session_state.get("active_background_job_id")
    if not job_id:
        st.caption("No active background batch job.")
        return

    job = get_background_job(job_id)
    if not job:
        st.caption("No active background batch job.")
        return

    job_status = job.get("status", "Idle")
    job_progress = int(job.get("progress", 0) or 0)
    job_total_files = int(job.get("total_files", 0) or 0)
    job_processed_files = int(job.get("processed_files", 0) or 0)
    job_current_file = job.get("current_file")
    job_results = job.get("results", []) or []
    job_exceptions = job.get("exceptions", []) or []
    job_zip_path = job.get("output_zip_path")
    job_zip_name = job.get("output_zip_name")
    job_pdf_path = job.get("assessment_pdf_path")
    job_rankings = job.get("rankings", []) or []
    job_is_running = bool(job.get("is_running", False))
    last_error = job.get("last_error")
    last_update = job.get("last_update")
    resumable = can_resume_background_job(job)

    completed_count = len([x for x in job_results if x.get("status") == "Completed"])
    review_count = len([x for x in job_results if x.get("status") == "Review Needed"])
    exception_count = len([x for x in job_results if x.get("status") == "Exception"])

    m1, m2, m3 = st.columns(3)
    with m1:
        st.metric("Completed", completed_count)
    with m2:
        st.metric("Review Needed", review_count)
    with m3:
        st.metric("Exceptions", exception_count)

    c1, c2 = st.columns(2)

    with c1:
        if st.button("Refresh Job Status", use_container_width=True, key=f"refresh_bg_job_{job_id}"):
            st.rerun()

    with c2:
        if st.button(
            "Resume Background Job",
            use_container_width=True,
            disabled=not resumable,
            key=f"resume_bg_job_{job_id}"
        ):
            resume_background_batch_job()
            st.rerun()

    if job_is_running or str(job_status).startswith(("Queued", "Running", "Processing", "Resuming")):
        st.info("Background batch is running.")
    else:
        st.caption("Background batch is not active.")

    st.markdown(f"**Status:** {job_status}")
    st.markdown(f"**Processed:** {job_processed_files} / {job_total_files}")
    st.markdown(f"**Current File:** {job_current_file or '-'}")
    st.progress(max(0, min(100, job_progress)))

    if last_update:
        st.caption(f"Last update heartbeat: {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(last_update))}")

    if last_error:
        st.warning(f"Last error: {last_error}")

    if job_results:
        rows = []
        for item in job_results:
            debug_info = item.get("debug_info") or {}
            rows.append({
                "File": item.get("file_name"),
                "Type": item.get("doc_type"),
                "Status": item.get("status"),
                "OCR": "Yes" if item.get("ocr_used") else "No",
                "Reason": item.get("exception_reason") or "-",
                "Detected Type": debug_info.get("detected_doc_type") or "-",
                "Final Type": debug_info.get("final_doc_type") or item.get("doc_type") or "-",
                "Fallback": "Yes" if debug_info.get("resume_fallback_used") else "No",
                "Extraction": debug_info.get("extraction_mode") or "-",
            })
        st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True, height=220)
        
        rejected_items = [x for x in job_results if x.get("status") == "Exception"]
        if rejected_items:
            with st.expander("Rejected File Debug", expanded=False):
                for item in rejected_items:
                    st.markdown(f"**File:** {item.get('file_name')}")
                    st.markdown(f"**Reason:** {item.get('exception_reason')}")
                    debug_info = item.get("debug_info") or {}
                    st.markdown(f"**Detected Type:** {debug_info.get('detected_doc_type') or '-'}")
                    st.markdown(f"**Extraction Mode:** {debug_info.get('extraction_mode') or '-'}")
                    preview = debug_info.get("text_preview") or ""
                    if preview:
                        st.text_area(
                            f"Preview - {item.get('file_name')}",
                            value=preview,
                            height=180,
                            key=f"debug_preview_{job_id}_{item.get('file_name')}"
                        )

    if job_exceptions:
        st.warning(f"{len(job_exceptions)} file(s) ended in exception.")

    d1, d2 = st.columns(2)

    with d1:
        if job_zip_path:
            zip_bytes = load_job_binary(job_zip_path)
            if zip_bytes:
                st.download_button(
                    "Download Background Job Resumes",
                    data=zip_bytes,
                    file_name=job_zip_name or "background_job_resumes.zip",
                    mime="application/zip",
                    use_container_width=True,
                    key=f"job_resume_zip_download_{job_id}"
                )

    with d2:
        if job_pdf_path:
            pdf_bytes = load_job_binary(job_pdf_path)
            if pdf_bytes:
                st.download_button(
                    "Download Background Assessment PDF",
                    data=pdf_bytes,
                    file_name="BackgroundDetailedAssessment.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key=f"job_assessment_pdf_download_{job_id}"
                )

    if job_rankings:
        with st.expander("Background Job JD Ranking", expanded=False):
            ranking_rows = []
            for item in job_rankings:
                ranking_rows.append({
                    "Rank": item.get("rank"),
                    "Candidate": item.get("candidate_name"),
                    "File": item.get("file_name"),
                    "Overall": item.get("overall_score"),
                    "Recommendation": item.get("recommendation"),
                })
            st.dataframe(pd.DataFrame(ranking_rows), use_container_width=True, hide_index=True)


def push_job_notification(message):
    if "job_notifications" not in st.session_state:
        st.session_state["job_notifications"] = []
    notifications = list(st.session_state.get("job_notifications", []))
    notifications.append(message)
    st.session_state["job_notifications"] = notifications
            
# ------------------------------
# REVIEW / ACTIONS
# ------------------------------
def render_validation_summary():
    validation = st.session_state.get("validation_result") or {}
    issues = validation.get("issues", [])
    warnings = validation.get("warnings", [])
    passed = validation.get("passed", False)

    st.markdown("#### Validation")
    if passed:
        st.success("Ready for approval")
    else:
        st.warning("Needs review before approval")

    for item in issues:
        st.caption(f"• {item}")
    for item in warnings:
        st.caption(f"• {item}")


def render_confidence_table():
    confidence = st.session_state.get("confidence_map") or {}
    if not confidence:
        return

    rows = [{"Field": k, "Confidence": v.get("label", "-"), "Reason": v.get("reason", "-")} for k, v in confidence.items()]
    st.markdown("#### Confidence")
    st.dataframe(pd.DataFrame(rows), use_container_width=True, height=220, hide_index=True)


def refresh_review_scores():
    data = st.session_state.get("review_data") or {}
    doc_type = st.session_state.get("doc_type") or "other"
    st.session_state.validation_result = validate_document_data(data, doc_type)
    st.session_state.confidence_map = build_confidence_map(data, doc_type)


def render_invoice_review_form():
    data = st.session_state.get("review_data") or {}
    with st.form("invoice_review_form"):
        c1, c2 = st.columns(2)
        vendor = c1.text_input("Vendor", value=str(data.get("vendor") or data.get("supplier") or ""))
        invoice_number = c2.text_input("Invoice Number", value=str(data.get("invoice_number") or data.get("invoice_no") or ""))
        c3, c4 = st.columns(2)
        invoice_date = c3.text_input("Invoice Date", value=str(data.get("invoice_date") or ""))
        due_date = c4.text_input("Due Date", value=str(data.get("due_date") or ""))
        c5, c6, c7, c8 = st.columns(4)
        currency = c5.text_input("Currency", value=str(data.get("currency") or ""))
        subtotal = c6.text_input("Subtotal", value=str(data.get("subtotal") or ""))
        tax = c7.text_input("Tax", value=str(data.get("tax") or ""))
        total = c8.text_input("Total", value=str(data.get("total") or ""))

        saved = st.form_submit_button("Save Review Changes", use_container_width=True)

    if saved:
        data["vendor"] = vendor
        data["invoice_number"] = invoice_number
        data["invoice_date"] = invoice_date
        data["due_date"] = due_date
        data["currency"] = currency
        data["subtotal"] = subtotal
        data["tax"] = tax
        data["total"] = total
        st.session_state.review_data = data
        refresh_review_scores()
        st.success("Review updates saved")


def render_ticket_review_form():
    data = st.session_state.get("review_data") or {}
    with st.form("ticket_review_form"):
        c1, c2 = st.columns(2)
        traveler_name = c1.text_input("Traveler Name", value=str(data.get("traveler_name") or ""))
        ticket_number = c2.text_input("Ticket Number", value=str(data.get("ticket_number") or ""))
        c3, c4 = st.columns(2)
        airline = c3.text_input("Airline", value=str(data.get("airline") or ""))
        booking_reference = c4.text_input("Booking Reference", value=str(data.get("booking_reference") or ""))
        c5, c6 = st.columns(2)
        from_city = c5.text_input("From", value=str(data.get("from") or ""))
        to_city = c6.text_input("To", value=str(data.get("to") or ""))
        c7, c8, c9 = st.columns(3)
        departure_date = c7.text_input("Departure Date", value=str(data.get("departure_date") or ""))
        return_date = c8.text_input("Return Date", value=str(data.get("return_date") or ""))
        amount = c9.text_input("Amount", value=str(data.get("amount") or ""))

        saved = st.form_submit_button("Save Review Changes", use_container_width=True)

    if saved:
        data["traveler_name"] = traveler_name
        data["ticket_number"] = ticket_number
        data["airline"] = airline
        data["booking_reference"] = booking_reference
        data["from"] = from_city
        data["to"] = to_city
        data["departure_date"] = departure_date
        data["return_date"] = return_date
        data["amount"] = amount
        st.session_state.review_data = data
        refresh_review_scores()
        st.success("Review updates saved")


def render_resume_review_form():
    data = st.session_state.get("review_data") or {}
    with st.form("resume_review_form"):
        c1, c2 = st.columns(2)
        name = c1.text_input("Name", value=str(data.get("name") or ""))
        email = c2.text_input("Email", value=str(data.get("email") or ""))
        c3, c4 = st.columns(2)
        phone = c3.text_input("Phone", value=str(data.get("phone") or ""))
        location = c4.text_input("Location", value=str(data.get("location") or ""))
        linkedin = st.text_input("LinkedIn", value=str(data.get("linkedin") or ""))
        skills = st.text_input("Skills (comma-separated)", value=", ".join(data.get("skills", [])) if isinstance(data.get("skills"), list) else "")
        summary = st.text_area("Summary", value=str(data.get("summary") or ""), height=120)
        saved = st.form_submit_button("Save Review Changes", use_container_width=True)

    if saved:
        data["name"] = name
        data["email"] = email
        data["phone"] = phone
        data["location"] = location
        data["linkedin"] = linkedin
        data["skills"] = [s.strip() for s in skills.split(",") if s.strip()]
        data["summary"] = summary
        st.session_state.review_data = data
        refresh_review_scores()
        st.success("Review updates saved")


def handle_invoice_or_ticket_submission(doc_type):
    validation = st.session_state.get("validation_result") or {}
    if not validation.get("passed"):
        st.warning("Please resolve validation issues before approval")
        return

    data = st.session_state.get("review_data") or {}
    result = send_to_concur(doc_type, data, mode="mock")
    st.session_state.auto_result["result"].update({
        "payload": result.get("payload"),
        "concur_status": result.get("status"),
        "concur_mode": result.get("mode"),
        "concur_submission_id": result.get("submission_id"),
        "concur_batch_id": result.get("batch_id"),
        "concur_document_id": result.get("document_id"),
        "concur_submitted_at": result.get("submitted_at"),
        "concur_endpoint": result.get("endpoint"),
        "concur_processing_state": result.get("processing_state"),
        "concur_next_status": result.get("next_status"),
        "message": result.get("message"),
    })

    save_version_snapshot(
        file_name=st.session_state.get("current_file"),
        doc_type=doc_type,
        review_data=st.session_state.get("review_data"),
        auto_result=st.session_state.get("auto_result"),
        status="Submitted",
        note=f"{doc_type.title()} submitted to Concur"
    )

    st.success(f"{doc_type.title()} approved and submitted to Concur")


def regenerate_resume_from_review():
    validation = st.session_state.get("validation_result") or {}
    data = st.session_state.get("review_data") or {}
    template_bytes = get_active_template_bytes()

    if not template_bytes:
        st.error("No resume template available")
        return

    if not validation.get("passed"):
        st.warning("Resume has validation issues. Review before regenerating.")
        return

    try:
        file_bytes = build_resume(data, template_bytes)
        st.session_state.generated_resume = file_bytes
        st.session_state.auto_result["result"]["file"] = file_bytes
        st.session_state.auto_result["result"]["data"] = data

        save_version_snapshot(
            file_name=st.session_state.get("current_file"),
            doc_type="resume",
            review_data=st.session_state.get("review_data"),
            auto_result=st.session_state.get("auto_result"),
            status="Regenerated",
            note="Resume regenerated after review edits"
        )

        st.success("Resume regenerated successfully")
    except Exception as e:
        st.error(f"Resume regeneration failed: {str(e)}")

# ------------------------------
# UI
# ------------------------------
def render_header():
    logo_path = Path(__file__).parent / "ResumeProcessor.png"
    col_logo, col_title = st.columns([1, 6], gap="small")

    with col_logo:
        if logo_path.exists():
            st.image(logo_path, width=300)

    with col_title:
        st.markdown("## Intelligent Resume Processor")
        st.caption("AI-powered resume evaluation & automation")

def render_sidebar_and_upload():
    with st.sidebar:
        st.write(f"Hi **{st.session_state['user']}**")
        st.markdown("---")

        model_choice = st.selectbox(
            "Choose Model",
            ["gpt-4o-mini", "gpt-4o", "gpt-5"],
            index=["gpt-4o-mini", "gpt-4o", "gpt-5"].index(
                st.session_state.get("model_choice", "gpt-4o-mini")
            )
        )
        st.session_state["model_choice"] = model_choice

        st.markdown("---")
        st.success("🔑 API key loaded securely")
        cost = st.session_state.get("metrics", {}).get("cost", 0.0)

        st.markdown("---")
        st.write(f"💰 Session Cost ${round(cost, 6)}")

        st.markdown("---")
        if st.button("Logout", use_container_width=True):
            for key in ["logged_in", "user", "role", "api_key"]:
                if key in st.session_state:
                    del st.session_state[key]
            st.rerun()

    st.markdown("### Document Source")

    source_options = ["Local Upload", "SharePoint", "OneDrive"]
    saved_source_mode = st.session_state.get("source_mode", "Local Upload")

    if saved_source_mode not in source_options:
        saved_source_mode = "Local Upload"
        st.session_state["source_mode"] = "Local Upload"

    source_mode = st.radio(
        "Choose source",
        source_options,
        horizontal=True,
        index=source_options.index(saved_source_mode),
    )

    st.session_state["source_mode"] = source_mode

    uploaded_files = []

    if source_mode == "Local Upload":
        c1, c2 = st.columns([6, 1], gap="small")

        with c1:
            uploaded_files = st.file_uploader(
                f"Upload resume document(s) - max {MAX_BATCH_FILES} files per batch",
                type=["txt", "pdf", "docx"],
                accept_multiple_files=True,
                key=f"main_file_uploader_{st.session_state.uploader_key}"
            ) or []

        with c2:
            st.write("")
            st.write("")
            if st.button("Reset", use_container_width=True):
                st.session_state.batch_results = []
                st.session_state.exception_queue = []
                st.session_state.active_batch_index = 0
                st.session_state.batch_processed = False
                st.session_state.last_batch_signature = None
                st.session_state.show_reprocess_confirm = False
                st.session_state.pending_batch_signature = None
                st.session_state.batch_total_files = 0
                st.session_state.batch_processed_files = 0
                st.session_state.batch_current_file = None
                st.session_state.batch_file_statuses = []
                st.session_state.jd_rankings = []
                st.session_state.jd_text = ""
                st.session_state.remote_uploaded_files = []
                st.session_state.source_mode = "Local Upload"
                st.session_state.uploader_key += 1
                reset_run_state()
                st.rerun()

    elif source_mode == "SharePoint":
        st.info("Load CVs from SharePoint using either site details or a folder URL.")
        st.caption("Supported remote CV types: PDF, DOCX, TXT")

        sp_mode = st.radio(
            "SharePoint input mode",
            ["Site + Library + Folder", "Folder URL"],
            horizontal=True,
            key="sp_input_mode"
        )

        c1, c2 = st.columns([6, 1], gap="small")

        if sp_mode == "Site + Library + Folder":
            s1, s2, s3 = st.columns(3)
            with s1:
                sharepoint_url = st.text_input(
                    "SharePoint Site URL",
                    placeholder="https://yourtenant.sharepoint.com/sites/Recruiting",
                    key="sp_site_url"
                )
            with s2:
                library_name = st.text_input(
                    "Library Name",
                    value="Documents",
                    key="sp_library_name"
                )
            with s3:
                folder_path = st.text_input(
                    "Folder Path",
                    value="CVs",
                    key="sp_folder_path"
                )

            with c1:
                if st.button("Load Resume CVs from SharePoint", use_container_width=True):
                    if not sharepoint_url or not folder_path:
                        st.warning("Please enter SharePoint Site URL and Folder Path.")
                    else:
                        try:
                            with st.spinner("Loading CVs from SharePoint..."):
                                fetched = get_cv_files_from_sharepoint(
                                    site_url=sharepoint_url,
                                    folder_path=folder_path,
                                    library_name=library_name or "Documents",
                                )
                                remote_files = [
                                    RemoteUploadedFile(name=f["name"], content=f["content"])
                                    for f in fetched
                                ]
                                st.session_state["remote_uploaded_files"] = remote_files
                                st.success(f"Loaded {len(remote_files)} CV file(s)")
                        except Exception as e:
                            st.error(f"Failed to load files: {str(e)}")
        else:
            folder_url = st.text_input(
                "SharePoint Folder URL",
                placeholder="Paste a SharePoint folder sharing URL",
                key="sp_folder_url"
            )

            with c1:
                if st.button("Load Resume CVs from SharePoint URL", use_container_width=True):
                    if not folder_url:
                        st.warning("Please paste a SharePoint folder URL.")
                    else:
                        try:
                            with st.spinner("Loading CVs from SharePoint URL..."):
                                fetched = get_cv_files_from_sharepoint_url(folder_url)
                                remote_files = [
                                    RemoteUploadedFile(name=f["name"], content=f["content"])
                                    for f in fetched
                                ]
                                st.session_state["remote_uploaded_files"] = remote_files
                                st.success(f"Loaded {len(remote_files)} CV file(s)")
                        except Exception as e:
                            st.error(f"Failed to load files: {str(e)}")

        remote_files = st.session_state.get("remote_uploaded_files", [])
        if remote_files:
            st.caption(f"{len(remote_files)} CV file(s) ready from SharePoint")
            st.write("Loaded files:")
            for name in [f.name for f in remote_files[:10]]:
                st.caption(f"• {name}")

        with c2:
            st.write("")
            st.write("")
            if st.button("Reset", use_container_width=True, key="reset_sharepoint_source"):
                reset_source_state()
                st.rerun()

        uploaded_files = st.session_state.get("remote_uploaded_files", [])

    else:
        st.info("Load CVs from OneDrive using either drive details or a shared folder URL.")
        st.caption("Supported remote CV types: PDF, DOCX, TXT")

        od_mode = st.radio(
            "OneDrive input mode",
            ["Drive ID + Folder Path", "Shared Folder URL"],
            horizontal=True,
            key="od_input_mode"
        )

        c1, c2 = st.columns([6, 1], gap="small")

        if od_mode == "Drive ID + Folder Path":
            s1, s2 = st.columns(2)
            with s1:
                drive_id = st.text_input(
                    "OneDrive Drive ID",
                    key="od_drive_id"
                )
            with s2:
                folder_path = st.text_input(
                    "Folder Path",
                    value="CVs",
                    key="od_folder_path"
                )

            with c1:
                if st.button("Load Resume CVs from OneDrive", use_container_width=True):
                    if not drive_id or not folder_path:
                        st.warning("Please enter OneDrive Drive ID and Folder Path.")
                    else:
                        try:
                            with st.spinner("Loading CVs from OneDrive..."):
                                fetched = get_cv_files_from_onedrive(
                                    drive_id=drive_id,
                                    folder_path=folder_path,
                                )
                                remote_files = [
                                    RemoteUploadedFile(name=f["name"], content=f["content"])
                                    for f in fetched
                                ]
                                st.session_state["remote_uploaded_files"] = remote_files
                                st.success(f"Loaded {len(remote_files)} CV file(s)")
                        except Exception as e:
                            st.error(f"Failed to load files: {str(e)}")
        else:
            shared_url = st.text_input(
                "OneDrive Shared Folder URL",
                placeholder="Paste a OneDrive shared folder URL",
                key="od_shared_url"
            )

            with c1:
                if st.button("Load Resume CVs from OneDrive URL", use_container_width=True):
                    if not shared_url:
                        st.warning("Please paste a OneDrive shared folder URL.")
                    else:
                        try:
                            with st.spinner("Loading CVs from OneDrive URL..."):
                                fetched = get_cv_files_from_onedrive_url(shared_url)
                                remote_files = [
                                    RemoteUploadedFile(name=f["name"], content=f["content"])
                                    for f in fetched
                                ]
                                st.session_state["remote_uploaded_files"] = remote_files
                                st.success(f"Loaded {len(remote_files)} CV file(s)")
                        except Exception as e:
                            st.error(f"Failed to load files: {str(e)}")

        remote_files = st.session_state.get("remote_uploaded_files", [])
        if remote_files:
            st.caption(f"{len(remote_files)} CV file(s) ready from OneDrive")
            st.write("Loaded files:")
            for name in [f.name for f in remote_files[:10]]:
                st.caption(f"• {name}")

        with c2:
            st.write("")
            st.write("")
            if st.button("Reset", use_container_width=True, key="reset_onedrive_source"):
                reset_source_state()
                st.rerun()

        uploaded_files = st.session_state.get("remote_uploaded_files", [])

    if uploaded_files and len(uploaded_files) > MAX_BATCH_FILES:
        st.error(f"Batch limit exceeded for direct processing. Maximum allowed is {MAX_BATCH_FILES} files.")

    st.markdown("---")
    return uploaded_files



def render_duplicate_warning():
    duplicate_info = st.session_state.get("duplicate_info") or {}
    if duplicate_info.get("is_duplicate"):
        st.warning(
            f"Possible duplicate detected with: {duplicate_info.get('match_file')} "
            f"({duplicate_info.get('reason')}, score={duplicate_info.get('score')})"
        )


def render_result_workspace():
    st.markdown("### Result")

    if not st.session_state.get("auto_result"):
        st.caption("Process a document to view results.")
        return

    doc_type = st.session_state.get("doc_type")
    result = st.session_state.get("auto_result", {}).get("result", {})
    data = st.session_state.get("review_data") or {}

    current_index = st.session_state.get("active_batch_index", 0)
    total_results = len(st.session_state.get("batch_results", []))
    has_next = current_index < (total_results - 1)

    if doc_type == "invoice":
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            compact_field("Vendor", str(data.get("vendor") or data.get("supplier") or "-"))
        with c2:
            compact_field("Invoice No", str(data.get("invoice_number") or data.get("invoice_no") or "-"))
        with c3:
            compact_field("Date", str(data.get("invoice_date") or "-"))
        with c4:
            compact_field("Total", str(data.get("total") or "-"))

        if st.session_state.get("auto_result", {}).get("ocr_used"):
            st.caption("OCR Applied")

        render_validation_summary()
        render_duplicate_warning()
        render_confidence_table()

        with st.expander("Review & Edit", expanded=True):
            render_invoice_review_form()

        b1, b2, b3, b4 = st.columns(4)
        with b1:
            if st.button("Approve & Send to Concur", use_container_width=True, key="invoice_send"):
                handle_invoice_or_ticket_submission("invoice")

        with b2:
            excel = result.get("excel")
            if excel:
                st.download_button(
                    "Download Excel",
                    excel,
                    f"{(data.get('invoice_number') or data.get('vendor') or 'invoice_data')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )

        with b3:
            json_bytes = json.dumps(data, indent=2, ensure_ascii=False).encode("utf-8")
            json_name = f"{(data.get('invoice_number') or data.get('vendor') or 'invoice_data')}.json"
            st.download_button(
                "Download JSON",
                json_bytes,
                json_name,
                mime="application/json",
                use_container_width=True,
                key="invoice_json_download"
            )

        with b4:
            if st.button("Next Document", use_container_width=True, disabled=not has_next, key="invoice_next"):
                go_to_next_batch_result()
                st.rerun()
                
    elif doc_type == "ticket":
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            compact_field("Traveler", str(data.get("traveler_name") or "-"))
        with c2:
            compact_field("Airline", str(data.get("airline") or "-"))
        with c3:
            compact_field("Route", f"{data.get('from', '-')}" + " → " + f"{data.get('to', '-')}")
        with c4:
            compact_field("Amount", str(data.get("amount") or "-"))

        if st.session_state.get("auto_result", {}).get("ocr_used"):
            st.caption("OCR Applied")

        render_validation_summary()
        render_duplicate_warning()
        render_confidence_table()

        with st.expander("Review & Edit", expanded=True):
            render_ticket_review_form()

        a1, a2, a3 = st.columns(3)
        with a1:
            if st.button("Approve & Send to Concur", use_container_width=True, key="ticket_send"):
                handle_invoice_or_ticket_submission("ticket")

        with a2:
            json_bytes = json.dumps(data, indent=2, ensure_ascii=False).encode("utf-8")
            json_name = f"{(data.get('ticket_number') or data.get('traveler_name') or 'ticket_data')}.json"
            st.download_button(
                "Download JSON",
                json_bytes,
                json_name,
                mime="application/json",
                use_container_width=True,
                key="ticket_json_download"
            )

        with a3:
            if st.button("Next Document", use_container_width=True, disabled=not has_next, key="ticket_next"):
                go_to_next_batch_result()
                st.rerun()
                
    elif doc_type == "resume":
        st.caption(f"Output File: {result.get('file_name', 'generated_resume.docx')}")

        r1, r2, r3 = st.columns(3)
        with r1:
            if st.button("Regenerate Resume", use_container_width=True, key="resume_regen"):
                regenerate_resume_from_review()

        with r2:
            if result.get("file"):
                st.download_button(
                    "Download Resume",
                    data=result["file"],
                    file_name=result.get("file_name", "generated_resume.docx"),
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )

        with r3:
            if st.button("Next Document", use_container_width=True, disabled=not has_next, key="resume_next"):
                go_to_next_batch_result()
                st.rerun()
                
        render_validation_summary()
        render_duplicate_warning()
        render_confidence_table()

        with st.expander("Review & Edit", expanded=True):
            render_resume_review_form()
        
    else:
        text = st.session_state.get("full_text", "")
        if text:
            st.text_area("Preview", value=text[:2500], height=180, label_visibility="collapsed")

        g1, g2 = st.columns(2)

        with g1:
            if st.button("Chat with Document", use_container_width=True, key="generic_chat"):
                st.session_state["open_doc_chat"] = True

        with g2:
            if st.button("Next Document", use_container_width=True, disabled=not has_next, key="generic_next"):
                go_to_next_batch_result()
                st.rerun()

        if st.session_state.get("open_doc_chat"):
            st.markdown("#### Document Chat")

            user_q = st.text_input("Ask a question about this document", key="generic_doc_chat_q")

            if st.button("Ask", use_container_width=True, key="generic_doc_chat_ask"):
                full_text = st.session_state.get("full_text", "")
                if not full_text.strip():
                    st.warning("No document text available for chat.")
                else:
                    try:
                        llm = get_llm(st.session_state["api_key"], st.session_state.get("model_choice", "gpt-4o-mini"))
                        prompt = f"""
Answer the user's question using only the document text below.
If the answer is not in the document, say so clearly.

DOCUMENT TEXT:
{full_text[:12000]}

QUESTION:
{user_q}
"""
                        answer = llm.invoke(prompt).content
                        st.session_state.setdefault("generic_doc_chat_history", []).append({
                            "question": user_q,
                            "answer": answer,
                        })
                    except Exception as e:
                        st.error(f"Chat failed: {str(e)}")

            history = st.session_state.get("generic_doc_chat_history", [])
            if history:
                for item in reversed(history[-5:]):
                    st.markdown(f"**Q:** {item['question']}")
                    st.markdown(f"**A:** {item['answer']}")

def render_batch_table():
    st.markdown("### Batch Results")
    elapsed = st.session_state.get("batch_elapsed_seconds", 0.0)
    if elapsed:
        st.caption(f"Batch processed in {elapsed:.2f} sec")

    batch_results = st.session_state.get("batch_results", [])
    if not batch_results:
        st.caption("No batch results yet")
        return

    rows = []
    for item in batch_results:
        dup = item.get("duplicate_info") or {}
        rows.append({
            "File": item.get("file_name"),
            "Type": item.get("doc_type"),
            "Status": item.get("status"),
            "OCR": "Yes" if item.get("ocr_used") else "No",
            "Duplicate": "Yes" if dup.get("is_duplicate") else "No",
            "Cost": item.get("cost"),
            "Tokens": item.get("tokens"),
        })

    df = pd.DataFrame(rows)
    st.dataframe(df, use_container_width=True, hide_index=True, height=220)

    current_index = st.session_state.get("active_batch_index", 0)
    if current_index < 0 or current_index >= len(batch_results):
        current_index = 0
        st.session_state["active_batch_index"] = 0

    selected = st.selectbox(
        "Open processed document",
        options=list(range(len(batch_results))),
        format_func=lambda i: f"{batch_results[i]['file_name']} ({batch_results[i]['status']})",
        index=current_index,
        key="batch_result_selector",
    )

    if selected is not None and 0 <= selected < len(batch_results):
        load_batch_result_into_session(selected)

def render_exception_queue():
    st.markdown("### Exception Queue")
    if not st.session_state.exception_queue:
        st.caption("No exceptions")
        return

    rows = []
    for item in st.session_state.exception_queue:
        rows.append({
            "File": item.get("file_name"),
            "Type": item.get("doc_type"),
            "Reason": item.get("exception_reason"),
            "OCR": "Yes" if item.get("ocr_used") else "No",
        })

    df = pd.DataFrame(rows)
    st.dataframe(df, use_container_width=True, hide_index=True, height=200)


def render_template_manager():
    st.markdown("### Template Manager")

    template_upload = st.file_uploader(
        "Upload Resume Template",
        type=["docx"],
        key="template_manager_uploader"
    )

    if template_upload and st.button("Add Template", use_container_width=True):
        add_template_to_library(template_upload)
        st.success("Template added to library")
        st.rerun()

    library = st.session_state.get("template_library", [])
    if not library:
        st.caption("No custom templates uploaded. Default template will be used.")
        return

    selected = st.selectbox(
        "Choose active template",
        options=list(range(len(library))),
        format_func=lambda i: library[i]["name"],
        index=st.session_state.get("active_template_index", 0) if library else 0,
        key="active_template_selector"
    )
    st.session_state.active_template_index = selected

    active = library[selected]
    validation = active.get("validation", {})

    if validation.get("valid"):
        st.success("Template is valid")
    else:
        st.warning("Template is missing required placeholders")

    with st.expander("Template Details", expanded=False):
        st.write("Found placeholders:")
        st.write(", ".join(validation.get("found_placeholders", [])) or "-")

        missing = validation.get("missing_placeholders", [])
        if missing:
            st.write("Missing placeholders:")
            st.write(", ".join(missing))


def render_version_history():
    st.markdown("### Version History")

    history = st.session_state.get("version_history", [])
    current_file = st.session_state.get("current_file")

    if not history:
        st.caption("No version history yet")
        return

    filtered = [h for h in history if h.get("file_name") == current_file] if current_file else history

    if not filtered:
        st.caption("No history for current document")
        return

    rows = []
    for idx, item in enumerate(filtered):
        rows.append({
            "Version": idx + 1,
            "Timestamp": item.get("timestamp"),
            "Status": item.get("status"),
            "Note": item.get("note"),
        })

    st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True, height=180)

    selected = st.selectbox(
        "Open version snapshot",
        options=list(range(len(filtered))),
        format_func=lambda i: f"Version {i+1} - {filtered[i]['timestamp']} - {filtered[i]['status']}",
        key="version_history_selector"
    )

    if selected is not None:
        snap = filtered[selected]
        with st.expander("Snapshot Details", expanded=False):
            st.json(snap.get("review_data", {}))


def render_batch_downloads():
    st.markdown("### Batch Downloads")

    resume_count, _ = get_batch_download_counts()

    if resume_count > 0:
        resume_zip = build_zip_from_batch_results("resume")
        st.download_button(
            label=f"Download All Resumes ({resume_count})",
            data=resume_zip,
            file_name="all_resumes.zip",
            mime="application/zip",
            use_container_width=True
        )
    else:
        st.button(
            "Download All Resumes (0)",
            disabled=True,
            use_container_width=True,
            key="dl_all_resumes_disabled"
        )


def render_jd_ranking():
    st.markdown("### JD Match Ranking")

    c1, c2 = st.columns([2, 1], gap="medium")

    with c1:
        pasted_jd = st.text_area(
            "Paste Job Description",
            value=st.session_state.get("jd_text", ""),
            height=220,
            key="jd_text_area"
        )

    with c2:
        jd_file = st.file_uploader(
            "Upload JD File",
            type=["pdf", "docx"],
            key="jd_file_uploader"
        )

        use_uploaded_jd = st.checkbox(
            "Use uploaded JD file",
            value=bool(jd_file),
            key="use_uploaded_jd_checkbox"
        )

    jd_text = pasted_jd.strip()

    if jd_file and use_uploaded_jd:
        uploaded_jd_text = extract_jd_text_from_upload(jd_file)
        if uploaded_jd_text:
            jd_text = uploaded_jd_text
            with st.expander("Preview extracted JD text", expanded=False):
                st.text_area(
                    "Extracted JD",
                    value=uploaded_jd_text[:5000],
                    height=200,
                    disabled=True,
                    label_visibility="collapsed"
                )

    st.session_state.jd_text = jd_text

    if st.button("Rank All CVs Against JD", use_container_width=True):
        rank_all_resumes_against_jd()

    rankings = st.session_state.get("jd_rankings", [])
    if not rankings:
        st.caption("No JD rankings yet")
        return

    rows = []
    for item in rankings:
        rows.append({
            "Rank": item.get("rank"),
            "Candidate": item.get("candidate_name"),
            "File": item.get("file_name"),
            "Overall": item.get("overall_score"),
            "Skills": item.get("skills_score"),
            "Experience": item.get("experience_score"),
            "Education": item.get("education_score"),
            "Recommendation": item.get("recommendation"),
        })

    st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True, height=240)

    selected = st.selectbox(
        "Open candidate analysis",
        options=list(range(len(rankings))),
        format_func=lambda i: f"#{rankings[i].get('rank')} - {rankings[i].get('candidate_name')} ({rankings[i].get('overall_score')})",
        key="jd_candidate_selector"
    )

    if selected is not None:
        item = rankings[selected]
        with st.expander("Candidate Analysis", expanded=True):
            st.markdown(f"**Candidate:** {item.get('candidate_name', '-')}")
            st.markdown(f"**Overall Score:** {item.get('overall_score', 0)}")
            st.markdown(f"**Recommendation:** {item.get('recommendation', '-')}")
            st.markdown("**Matched Skills**")
            st.write(", ".join(item.get("matched_skills", [])) or "-")
            st.markdown("**Missing Skills**")
            st.write(", ".join(item.get("missing_skills", [])) or "-")
            st.markdown("**Strengths**")
            for s in item.get("strengths", []):
                st.caption(f"• {s}")
            st.markdown("**Gaps**")
            for g in item.get("gaps", []):
                st.caption(f"• {g}")

def render_detailed_assessment_report():
    st.markdown("### Detailed Assessment Report")

    if st.button("Generate Detailed Assessment", use_container_width=True):
        jd_text = (st.session_state.get("jd_text") or "").strip()
        batch_results = st.session_state.get("batch_results", [])
        jd_rankings = st.session_state.get("jd_rankings", [])

        resume_count = len([
            x for x in batch_results
            if x.get("doc_type") == "resume" and x.get("review_data")
        ])

        if resume_count == 0:
            st.warning("No processed resumes available in the batch.")
            return

        if not jd_text:
            st.warning("Please provide a JD first using paste or upload.")
            return

        if not jd_rankings:
            st.warning("Please run JD ranking first.")
            return

        report_data = generate_consolidated_assessment_data(
            batch_results=batch_results,
            jd_text=jd_text,
            jd_rankings=jd_rankings
        )
        pdf_bytes = build_consolidated_assessment_pdf(report_data)

        st.session_state["detailed_assessment_data"] = report_data
        st.session_state["detailed_assessment_pdf"] = pdf_bytes
        st.success("Detailed assessment generated successfully.")

    report_data = st.session_state.get("detailed_assessment_data")
    pdf_bytes = st.session_state.get("detailed_assessment_pdf")

    if not report_data:
        st.caption("No detailed assessment generated yet.")
        return

    executive = report_data.get("executive_summary", {})
    candidates = report_data.get("candidates", [])
    final_summary = report_data.get("final_summary", {})
    recruiter_questions = report_data.get("recruiter_questions", [])

    k1, k2, k3 = st.columns(3)
    with k1:
        st.metric("Candidates", executive.get("total_candidates", 0))
    with k2:
        st.metric("Top Match Range", executive.get("top_match_range", "-"))
    with k3:
        st.metric("Recommended Action", executive.get("recommended_action", "-"))

    with st.expander("Executive Summary", expanded=True):
        st.markdown("**JD Summary**")
        st.write(executive.get("jd_summary", "-"))
        st.markdown("**Executive Takeaway**")
        st.write(executive.get("executive_takeaway", "-"))

    st.markdown("#### Candidate Score Cards")
    if candidates:
        score_cols = st.columns(min(4, len(candidates)))
        for idx, candidate in enumerate(candidates[:4]):
            with score_cols[idx]:
                score = int(candidate.get("overall_score", 0))
                if score >= 85:
                    box_color = "#e7f8ee"
                    text_color = "#0f9d58"
                elif score >= 70:
                    box_color = "#fff4e5"
                    text_color = "#b26a00"
                else:
                    box_color = "#fdecec"
                    text_color = "#b42318"

                st.markdown(
                    f"""
                    <div style="
                        border-radius:16px;
                        padding:16px;
                        background:{box_color};
                        border:1px solid #e5e7eb;
                        min-height:130px;
                    ">
                        <div style="font-size:14px;font-weight:700;color:#111827;">
                            {candidate.get('candidate_name', '-')}
                        </div>
                        <div style="font-size:28px;font-weight:800;color:{text_color};margin-top:6px;">
                            {score}
                        </div>
                        <div style="font-size:12px;color:#4b5563;">
                            {candidate.get('shortlist_label', '-')} • {candidate.get('recommendation', '-')}
                        </div>
                        <div style="font-size:12px;color:#6b7280;margin-top:8px;">
                            {candidate.get('current_role', '-') or '-'}
                        </div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )

    st.markdown("#### Candidate Ranking & Shortlist Decision")
    rows = []
    for idx, candidate in enumerate(candidates, start=1):
        rows.append({
            "Rank": idx,
            "Candidate": candidate.get("candidate_name"),
            "File": candidate.get("file_name"),
            "Overall Score": candidate.get("overall_score"),
            "Recommendation": candidate.get("shortlist_label"),
            "Fitment Progress": candidate.get("fitment_progress"),
        })
    if rows:
        st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

    st.markdown("#### Shortlist Recommendation")
    s1, s2, s3 = st.columns(3)
    with s1:
        st.success("Primary")
        for item in final_summary.get("primary_candidates", []):
            st.caption(f"• {item}")
    with s2:
        st.info("Backup")
        for item in final_summary.get("backup_candidates", []):
            st.caption(f"• {item}")
    with s3:
        st.warning("Hold")
        for item in final_summary.get("hold_candidates", []):
            st.caption(f"• {item}")

    if candidates:
        candidate_index = st.selectbox(
            "Open candidate detailed view",
            options=list(range(len(candidates))),
            format_func=lambda i: f"{i+1}. {candidates[i].get('candidate_name', '-')}",
            key="detailed_assessment_candidate_selector"
        )

        candidate = candidates[candidate_index]
        with st.expander("Candidate Detailed View", expanded=True):
            d1, d2, d3, d4 = st.columns(4)
            with d1:
                st.metric("Overall", candidate.get("overall_score", 0))
            with d2:
                st.metric("Skills", candidate.get("skills_score", 0))
            with d3:
                st.metric("Experience", candidate.get("experience_score", 0))
            with d4:
                st.metric("Education", candidate.get("education_score", 0))

            st.markdown(f"**Candidate:** {candidate.get('candidate_name', '-')}")
            st.markdown(f"**Current Role:** {candidate.get('current_role', '-') or '-'}")
            st.markdown(f"**Location:** {candidate.get('location', '-') or '-'}")
            st.markdown(f"**Recommendation:** {candidate.get('recommendation', '-')}")

            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**Matched Skills**")
                st.write(", ".join(candidate.get("matched_skills", [])) or "-")
                st.markdown("**Strengths**")
                for s in candidate.get("strengths", []):
                    st.caption(f"• {s}")
            with c2:
                st.markdown("**Missing Skills**")
                st.write(", ".join(candidate.get("missing_skills", [])) or "-")
                st.markdown("**Gaps / Risks**")
                for g in candidate.get("gaps", []):
                    st.caption(f"• {g}")

    with st.expander("Recruiter Screening Questions", expanded=False):
        if recruiter_questions:
            q_rows = []
            for item in recruiter_questions:
                q_rows.append({
                    "Question": item.get("question", "-"),
                    "Expected Answer": item.get("expected_answer", "-")
                })
            st.dataframe(pd.DataFrame(q_rows), use_container_width=True, hide_index=True)

    if pdf_bytes:
        st.download_button(
            "Download DetailedAssesment.pdf",
            data=pdf_bytes,
            file_name="DetailedAssesment.pdf",
            mime="application/pdf",
            use_container_width=True
        )

# ------------------------------
# MAIN
# ------------------------------
render_header()
uploaded_files = render_sidebar_and_upload()

left_col, right_col = st.columns([1, 1.6], gap="large")

with left_col:
    st.markdown("### Activity")
    st.session_state["live_pipeline_placeholder"] = st.empty()

    st.markdown("---")
    st.session_state["live_step_placeholder"] = st.empty()
    st.session_state["live_progress_placeholder"] = st.empty()
    st.session_state["live_event_placeholder"] = st.empty()    

    refresh_live_batch_activity()

    current_batch_signature = get_batch_signature(uploaded_files)
    last_batch_signature = st.session_state.get("last_batch_signature")

    process_disabled = not uploaded_files

    if st.button("Submit Background Job", use_container_width=True, disabled=process_disabled, key="submit_bg_job_btn"):
        submit_background_batch_job(uploaded_files)
    
    if st.button("Process Batch", use_container_width=True, disabled=process_disabled):
        if uploaded_files and len(uploaded_files) > MAX_BATCH_FILES:
            st.error(f"Direct processing is limited to {MAX_BATCH_FILES} files. Use 'Submit Background Job' for large batches.")
        else:
            if current_batch_signature and current_batch_signature == last_batch_signature:
                st.session_state.show_reprocess_confirm = True
                st.session_state.pending_batch_signature = current_batch_signature
            else:
                st.session_state.batch_results = []
                st.session_state.exception_queue = []
                st.session_state.jd_rankings = []
                st.session_state.show_reprocess_confirm = False
                st.session_state.pending_batch_signature = None
    
                st.session_state.batch_started_at = time.time()
                st.session_state.batch_completed_at = None
                st.session_state.batch_elapsed_seconds = 0.0
    
                st.session_state.batch_total_files = len(uploaded_files)
                st.session_state.batch_processed_files = 0
                st.session_state.batch_current_file = None
                st.session_state.batch_file_statuses = [
                    {"file_name": f.name, "status": "pending", "message": ""}
                    for f in uploaded_files
                ]
                refresh_live_batch_activity()
    
                for uploaded_file in uploaded_files:
                    try:
                        st.session_state["current_file_started_at"] = time.time()
                        st.session_state.batch_current_file = uploaded_file.name
                        update_batch_file_status(uploaded_file.name, "running", "Processing started")
                        refresh_live_batch_activity()
    
                        result = process_single_file(uploaded_file)
                        st.session_state.batch_results.append(result)
    
                        if result.get("status") == "Exception":
                            st.session_state.exception_queue.append(result)
                            update_batch_file_status(
                                uploaded_file.name,
                                "error",
                                result.get("exception_reason", "Exception")
                            )
                        elif result.get("status") == "Review Needed":
                            update_batch_file_status(uploaded_file.name, "done", "Review Needed")
                        else:
                            update_batch_file_status(
                                uploaded_file.name,
                                "done",
                                result.get("status", "Completed")
                            )
    
                    except Exception as e:
                        error_result = {
                            "file_name": uploaded_file.name,
                            "status": "Exception",
                            "doc_type": "unknown",
                            "ocr_used": False,
                            "exception_reason": f"Unhandled error: {str(e)}",
                            "cost": 0.0,
                            "tokens": 0,
                            "duplicate_info": {
                                "is_duplicate": False,
                                "match_file": None,
                                "reason": None,
                                "score": 0.0,
                            },
                            "agent_events": deepcopy(st.session_state.get("agent_events", [])),
                            "agent_timings": deepcopy(st.session_state.get("agent_timings", {})),
                        }
                        st.session_state.batch_results.append(error_result)
                        st.session_state.exception_queue.append(error_result)
                        update_batch_file_status(uploaded_file.name, "error", f"Unhandled error: {str(e)}")
    
                    finally:
                        st.session_state.batch_processed_files += 1
                        st.session_state["progress_value"] = 0
                        st.session_state["current_file_started_at"] = None
                        refresh_live_batch_activity()
    
                if st.session_state.batch_results:
                    load_batch_result_into_session(0)
                    st.session_state.batch_processed = True
                    st.session_state.last_batch_signature = current_batch_signature
                    st.session_state.batch_completed_at = time.time()
                    st.session_state.batch_elapsed_seconds = (
                        st.session_state.batch_completed_at - st.session_state.batch_started_at
                    )
                    st.success("Batch processing completed")
    
    if st.session_state.get("show_reprocess_confirm"):
        st.warning("This same batch was already processed. Do you want to re-process it again?")

        c1, c2 = st.columns(2)

        with c1:
            if st.button("Yes, Re-process", use_container_width=True):
                st.session_state.batch_results = []
                st.session_state.exception_queue = []
                st.session_state.jd_rankings = []
                st.session_state.batch_started_at = time.time()
                st.session_state.batch_completed_at = None
                st.session_state.batch_elapsed_seconds = 0.0

                st.session_state.batch_total_files = len(uploaded_files or [])
                st.session_state.batch_processed_files = 0
                st.session_state.batch_current_file = None
                st.session_state.batch_file_statuses = [
                    {"file_name": f.name, "status": "pending", "message": ""}
                    for f in (uploaded_files or [])
                ]
                refresh_live_batch_activity()

                for uploaded_file in (uploaded_files or []):
                    try:
                        st.session_state["current_file_started_at"] = time.time()
                        st.session_state.batch_current_file = uploaded_file.name
                        update_batch_file_status(uploaded_file.name, "running", "Re-processing started")
                        refresh_live_batch_activity()

                        result = process_single_file(uploaded_file)
                        st.session_state.batch_results.append(result)

                        if result.get("status") == "Exception":
                            st.session_state.exception_queue.append(result)
                            update_batch_file_status(
                                uploaded_file.name,
                                "error",
                                result.get("exception_reason", "Exception")
                            )
                        elif result.get("status") == "Review Needed":
                            update_batch_file_status(uploaded_file.name, "done", "Review Needed")
                        else:
                            update_batch_file_status(
                                uploaded_file.name,
                                "done",
                                result.get("status", "Completed")
                            )

                    except Exception as e:
                        error_result = {
                            "file_name": uploaded_file.name,
                            "status": "Exception",
                            "doc_type": "unknown",
                            "ocr_used": False,
                            "exception_reason": f"Unhandled error: {str(e)}",
                            "cost": 0.0,
                            "tokens": 0,
                            "duplicate_info": {
                                "is_duplicate": False,
                                "match_file": None,
                                "reason": None,
                                "score": 0.0,
                            },
                            "agent_events": deepcopy(st.session_state.get("agent_events", [])),
                            "agent_timings": deepcopy(st.session_state.get("agent_timings", {})),
                        }
                        st.session_state.batch_results.append(error_result)
                        st.session_state.exception_queue.append(error_result)
                        update_batch_file_status(uploaded_file.name, "error", f"Unhandled error: {str(e)}")

                    finally:
                        st.session_state.batch_processed_files += 1
                        st.session_state["progress_value"] = 0
                        st.session_state["current_file_started_at"] = None
                        refresh_live_batch_activity()

                if st.session_state.batch_results:
                    load_batch_result_into_session(0)
                    st.session_state.batch_processed = True
                    st.session_state.last_batch_signature = st.session_state.get("pending_batch_signature")
                    st.session_state.batch_completed_at = time.time()
                    st.session_state.batch_elapsed_seconds = (
                        st.session_state.batch_completed_at - st.session_state.batch_started_at
                    )
                    st.success("Batch re-processing completed")
                    
                st.session_state.show_reprocess_confirm = False
                st.session_state.pending_batch_signature = None
                st.rerun()

        with c2:
            if st.button("No", use_container_width=True):
                st.session_state.show_reprocess_confirm = False
                st.session_state.pending_batch_signature = None
                st.info("Re-processing cancelled")
                st.rerun()

with right_col:
    render_result_workspace()

st.markdown("---")
render_batch_table()
render_exception_queue()
render_batch_downloads()

st.markdown("---")
render_jd_ranking()

st.markdown("---")
render_detailed_assessment_report()

st.markdown("---")
render_template_manager()

st.markdown("---")
render_job_notifications()
render_background_job_monitor()

with st.expander("Metrics", expanded=False):
    m = st.session_state.get("metrics", {})
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Cost", f"${m.get('cost', 0.0):.6f}")
    c2.metric("Total Tokens", m.get("tokens", 0))
    c3.metric("Input Tokens", m.get("input_tokens", 0))
    c4.metric("Output Tokens", m.get("output_tokens", 0))

    doc_costs = st.session_state.get("doc_costs", {})
    doc_rows = [
        {"Document": k, "Cost": round(v.get("cost", 0.0), 6), "Tokens": v.get("tokens", 0)}
        for k, v in doc_costs.items()
    ]
    if doc_rows:
        st.dataframe(pd.DataFrame(doc_rows), use_container_width=True, hide_index=True, height=220)
